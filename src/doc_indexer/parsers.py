"""
Document parsers for different file formats.
"""

from pathlib import Path
from typing import Dict, Protocol

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from .models import Document, DocumentMetadata


class FileParser(Protocol):
    """Protocol for file parsers."""

    def parse(self, file_path: Path) -> Document:
        """Parse a file and return a Document."""
        ...


class PDFParser:
    """Parser for PDF files."""

    def parse(self, file_path: Path) -> Document:
        """Parse a PDF file."""
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        content_parts = []

        with open(file_path, "rb") as file:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    content_parts.append(text)

        content = "\n".join(content_parts)

        metadata = DocumentMetadata(
            filename=file_path.name,
            file_type="pdf",
            file_path=str(file_path.absolute()),
            file_size=file_path.stat().st_size,
        )

        return Document(content=content, metadata=metadata)


class WordParser:
    """Parser for Word documents."""

    def parse(self, file_path: Path) -> Document:
        """Parse a Word document."""
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        doc = DocxDocument(str(file_path))
        content_parts = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content_parts.append(paragraph.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        content_parts.append(cell.text)

        content = "\n".join(content_parts)

        metadata = DocumentMetadata(
            filename=file_path.name,
            file_type="docx",
            file_path=str(file_path.absolute()),
            file_size=file_path.stat().st_size,
        )

        return Document(content=content, metadata=metadata)


class PowerPointParser:
    """Parser for PowerPoint presentations."""

    def parse(self, file_path: Path) -> Document:
        """Parse a PowerPoint presentation."""
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        prs = Presentation(str(file_path))
        content_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_texts.append(text)

            if slide_texts:
                content_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_texts))

        content = "\n\n".join(content_parts)

        metadata = DocumentMetadata(
            filename=file_path.name,
            file_type="pptx",
            file_path=str(file_path.absolute()),
            file_size=file_path.stat().st_size,
        )

        return Document(content=content, metadata=metadata)


class DocumentParser:
    """Main document parser that delegates to specific parsers."""

    def __init__(self) -> None:
        """Initialize parser with supported file types."""
        self.parsers: Dict[str, FileParser] = {
            ".pdf": PDFParser(),
            ".docx": WordParser(),
            ".doc": WordParser(),  # Assuming .doc files can be handled similarly
            ".pptx": PowerPointParser(),
            ".ppt": PowerPointParser(),  # Assuming .ppt files can be handled similarly
        }

    def parse(self, file_path: Path) -> Document:
        """Parse a document based on its file extension."""
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_extension = file_path.suffix.lower()

        if file_extension not in self.parsers:
            raise ValueError(f"Unsupported file type: {file_extension}")

        parser = self.parsers[file_extension]
        return parser.parse(file_path)

    def is_supported(self, file_path: Path) -> bool:
        """Check if a file type is supported."""
        return file_path.suffix.lower() in self.parsers
