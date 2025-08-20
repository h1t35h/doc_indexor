"""
PowerPoint presentation parser with LLM support.
"""

import io
from pathlib import Path
from typing import List, Optional

from PIL import Image
from pptx import Presentation

from doc_indexer.parsers.base import BaseParser, PageContent


class PowerPointParser(BaseParser):
    """Parser for PowerPoint presentations with slide image extraction."""

    def __init__(self, *args, extract_images: bool = True, **kwargs) -> None:
        """
        Initialize PowerPoint parser.

        Args:
            extract_images: Whether to extract slide images
            *args, **kwargs: Arguments passed to BaseParser
        """
        super().__init__(*args, **kwargs)
        self.extract_images = extract_images

    async def extract_pages(self, file_path: Path) -> List[PageContent]:
        """Extract slides from PowerPoint presentation."""
        prs = Presentation(str(file_path))
        pages = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = self._extract_slide_text(slide)
            slide_tables = self._extract_slide_tables(slide)

            slide_image = None
            if self.extract_images:
                slide_image = self._extract_slide_image(slide, slide_num)

            page = PageContent(
                page_number=slide_num,
                text=slide_text,
                image=slide_image,
                tables=slide_tables,
            )
            pages.append(page)

        return pages

    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a slide."""
        text_parts = []

        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                text_parts.append(f"Title: {title_text}")

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                if shape == slide.shapes.title:
                    continue

                text = shape.text_frame.text.strip()
                if text:
                    if self._is_bulleted_text(shape.text_frame):
                        text_parts.append(self._format_bulleted_text(shape.text_frame))
                    else:
                        text_parts.append(text)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                text_parts.append(f"\nNotes: {notes_text}")

        return "\n\n".join(text_parts)

    def _is_bulleted_text(self, text_frame) -> bool:
        """
        Check if text frame contains bulleted text.

        Args:
            text_frame: PowerPoint text frame

        Returns:
            True if bulleted
        """
        for paragraph in text_frame.paragraphs:
            if paragraph.level > 0 or paragraph.text.strip().startswith(
                ("•", "-", "*")
            ):
                return True
        return False

    def _format_bulleted_text(self, text_frame) -> str:
        """
        Format bulleted text with proper indentation.

        Args:
            text_frame: PowerPoint text frame

        Returns:
            Formatted text
        """
        lines = []
        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                indent = "  " * paragraph.level
                if not text.startswith(("•", "-", "*")):
                    text = f"• {text}"
                lines.append(f"{indent}{text}")
        return "\n".join(lines)

    def _extract_slide_tables(self, slide) -> List[dict]:
        """
        Extract tables from a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            List of table data dictionaries
        """
        tables = []

        for shape in slide.shapes:
            if shape.has_table:
                table_data = {"header": [], "rows": []}

                table = shape.table
                for i, row in enumerate(table.rows):
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())

                    if i == 0:
                        # Assume first row is header
                        table_data["header"] = row_data
                    else:
                        table_data["rows"].append(row_data)

                tables.append(table_data)

        return tables

    def _extract_slide_image(self, slide, slide_num: int) -> Optional[Image.Image]:
        """
        Extract or generate an image of the slide.

        This is a simplified version. For full slide rendering,
        you would need to use python-pptx2img or similar library.

        Args:
            slide: PowerPoint slide object
            slide_num: Slide number

        Returns:
            PIL Image of the slide or None
        """
        # First, try to extract background image
        if hasattr(slide, "background") and slide.background:
            try:
                if slide.background.fill.type == 6:  # Picture fill
                    image_part = slide.background.fill.picture
                    if image_part:
                        image_data = image_part.blob
                        return Image.open(io.BytesIO(image_data))
            except Exception:
                pass

        # Extract embedded images from shapes
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture shape
                try:
                    image_part = shape.image
                    image_data = image_part.blob
                    image = Image.open(io.BytesIO(image_data))

                    # Resize if too large
                    if image.width > 1920 or image.height > 1080:
                        image.thumbnail((1920, 1080), Image.Resampling.LANCZOS)

                    return image  # Return first image found
                except Exception:
                    continue

        # Note: For complete slide rendering, you would need to:
        # 1. Use python-pptx2img to convert slide to image
        # 2. Or export slide as image using COM automation on Windows
        # 3. Or use a service like LibreOffice in headless mode

        return None

    def _extract_charts(self, slide) -> List[dict]:
        """
        Extract chart data from slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            List of chart data dictionaries
        """
        charts = []

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                chart_data = {
                    "type": str(chart.chart_type),
                    "title": chart.chart_title.text if chart.has_title else "",
                    "categories": [],
                    "series": [],
                }

                # Extract categories
                try:
                    for category in chart.plots[0].categories:
                        chart_data["categories"].append(str(category))
                except Exception:
                    pass

                # Extract series data
                try:
                    for series in chart.series:
                        series_data = {
                            "name": series.name,
                            "values": list(series.values),
                        }
                        chart_data["series"].append(series_data)
                except Exception:
                    pass

                charts.append(chart_data)

        return charts
